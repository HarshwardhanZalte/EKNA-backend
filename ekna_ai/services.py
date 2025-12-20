import io
import logging
import base64
import os
import tempfile

# Text Extractors
import fitz
import docx
import pptx
import pandas as pd

# Django & Models
from django.conf import settings
from ekna_app.models import Document
from ekna_ai.models import DocumentEmbedding
from ekna_app.utils import get_s3_client

# LangChain
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings

# Groq
from groq import Groq

# Supported File Extensions
SUPPORTED_FILE_EXTENSIONS = ['pdf', 'docx', 'doc', 'xlsx', 'xls', 'csv', 'pptx', 'ppt', 'jpg', 'jpeg', 'png', 'tiff', 'bmp', 'txt']

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        # LOCAL LIGHTWEIGHT MODEL
        self.embeddings_model = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2"
        )

        self.groq_client = Groq(api_key=settings.GROQ_API_KEY)

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=100
        )

        self.s3_client = get_s3_client()

    def process_document(self, doc_id: int):
        temp_file_path = None
        try:
            doc = Document.objects.get(id=doc_id)
            bucket_name = settings.AWS_S3_BUCKET_NAME
            file_ext = doc.s3_key.split('.')[-1].lower()

            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp_file:
                self.s3_client.download_fileobj(bucket_name, doc.s3_key, tmp_file)
                temp_file_path = tmp_file.name

            text_content = self._extract_text(temp_file_path, file_ext)

            if not text_content.strip():
                logger.warning(f"No text extracted from {doc.doc_name}")
                return

            DocumentEmbedding.objects.filter(doc=doc).delete()

            chunks = self.text_splitter.split_text(text_content)
            logger.info(f"Embedding {len(chunks)} chunks")

            BATCH_SIZE = 32

            for i in range(0, len(chunks), BATCH_SIZE):
                batch_chunks = chunks[i:i + BATCH_SIZE]

                vectors = self.embeddings_model.embed_documents(batch_chunks)

                objs = [
                    DocumentEmbedding(
                        doc=doc,
                        chunk=chunk,
                        chunk_index=i + idx,
                        embedding=vector  # 384-d
                    )
                    for idx, (chunk, vector) in enumerate(zip(batch_chunks, vectors))
                ]

                DocumentEmbedding.objects.bulk_create(objs)

            doc.is_processed = True
            doc.save()
            logger.info(f"Document processed: {doc.doc_name}")

        except Exception as e:
            logger.error(f"Document processing failed: {e}")

        finally:
            if temp_file_path and os.path.exists(temp_file_path):
                os.remove(temp_file_path)


    def _extract_text(self, file_path, ext) -> str:
        """
        Extracts text reading from the file path on disk.
        """
        text = []
        try:
            if ext == 'pdf':
                doc = fitz.open(file_path)
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    page_text = page.get_text().strip()
                    
                    # Check if it's a scanned PDF
                    if len(page_text) < 50: 
                        logger.info(f"Page {page_num+1} looks like a scan. Using Groq Vision...")
                        
                        pix = page.get_pixmap()
                        image_bytes = pix.tobytes("png")
                        
                        vision_text = self._extract_text_from_image_groq(image_bytes, "png")
                        text.append(vision_text)
                    else:
                        text.append(page_text)
            
            elif ext in ['docx', 'doc']:
                doc = docx.Document(file_path)
                text.append("\n".join([para.text for para in doc.paragraphs]))
                
            elif ext in ['xlsx', 'xls', 'csv']:
                if ext == 'csv':
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                text.append(df.to_string(index=False))
                
            elif ext in ['pptx', 'ppt']:
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text + "\n")
                            
            elif ext in ['jpg', 'jpeg', 'png', 'tiff', 'bmp']:
                # Read bytes from disk only for the image processing
                with open(file_path, "rb") as image_file:
                    image_bytes = image_file.read()
                    extracted = self._extract_text_from_image_groq(image_bytes, ext)
                    if extracted:
                        text.append(extracted)
                
            elif ext == 'txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text.append(f.read())

        except Exception as e:
            logger.error(f"Failed to extract text from {ext} file: {e}")
            
        return "\n".join(text)

    def _extract_text_from_image_groq(self, image_bytes, ext) -> str:
        """
        Uses llama-4-scout-17b-16e-instruct on Groq to extract text.
        """
        try:
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            mime_type = f"image/{ext}" if ext != 'jpg' else "image/jpeg"

            chat_completion = self.groq_client.chat.completions.create(
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Extract all the text present in this image. Output ONLY the raw text found, no conversational filler."},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{base64_image}",
                                },
                            },
                        ],
                    }
                ],
                model="meta-llama/llama-4-scout-17b-16e-instruct",
                temperature=0.0,
            )
            
            text = chat_completion.choices[0].message.content
            
            return text or ""
        except Exception as e:
            logger.error(f"Groq Vision Error: {e}")
            return ""